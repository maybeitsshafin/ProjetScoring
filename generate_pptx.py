"""
Génère une présentation PowerPoint complète du projet Scoring de Crédit
à partir des résultats des notebooks d'exploration et de modélisation.
"""

import json
import base64
import io
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Couleurs ──────────────────────────────────────────────
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
MEDIUM_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_BLUE = RGBColor(0x3A, 0x86, 0xC8)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)
MEDIUM_GRAY = RGBColor(0x7F, 0x8C, 0x8D)

ROOT = Path(__file__).parent
OUTPUT_DIR = ROOT / "outputs" / "figures" / "pptx_images"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def extract_images_from_notebook(nb_path: Path) -> dict[int, Path]:
    """Extrait les images PNG des cellules d'un notebook."""
    with open(nb_path, encoding="utf-8") as f:
        nb = json.load(f)

    images = {}
    for i, cell in enumerate(nb["cells"]):
        for output in cell.get("outputs", []):
            data = output.get("data", {})
            if "image/png" in data:
                img_data = base64.b64decode(data["image/png"])
                img_path = OUTPUT_DIR / f"{nb_path.stem}_cell{i}.png"
                img_path.write_bytes(img_data)
                images[i] = img_path
                break
    return images


def add_background(slide, color=DARK_BLUE):
    """Applique un fond de couleur à un slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """Ajoute un rectangle de couleur."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def set_text(tf, text, size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    """Configure le texte d'un text frame."""
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment


def add_text_box(slide, left, top, width, height, text, size=18, color=DARK_GRAY,
                 bold=False, alignment=PP_ALIGN.LEFT):
    """Ajoute une zone de texte."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    set_text(tf, text, size, color, bold, alignment)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, size=16, color=DARK_GRAY,
                    bullet_char="\u2022"):
    """Ajoute une liste à puces."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox


def add_metric_card(slide, left, top, width, height, label, value, color=LIGHT_BLUE):
    """Ajoute une carte métrique stylée."""
    shape = add_shape_bg(slide, left, top, width, height, color)

    # Valeur
    add_text_box(slide, left, top + Inches(0.15), width, Inches(0.5),
                 str(value), size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left, top + Inches(0.6), width, Inches(0.4),
                 label, size=12, color=WHITE, alignment=PP_ALIGN.CENTER)


def add_table_slide(slide, left, top, rows_data, col_widths, header_color=MEDIUM_BLUE):
    """Ajoute un tableau formaté."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                          sum(col_widths), Inches(0.4) * n_rows)
    table = table_shape.table

    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w

    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.alignment = PP_ALIGN.CENTER
                if ri == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return table_shape


def slide_title_bar(slide, title, subtitle=None):
    """Ajoute une barre de titre en haut du slide."""
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
                 title, size=28, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.65), Inches(12), Inches(0.4),
                     subtitle, size=14, color=RGBColor(0xBD, 0xC3, 0xC7))


# ══════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════
def main():
    print("Extraction des images des notebooks...")
    imgs_nb1 = extract_images_from_notebook(ROOT / "notebooks" / "01_exploration_nettoyage.ipynb")
    imgs_nb2 = extract_images_from_notebook(ROOT / "notebooks" / "02_feature_engineering_modeling.ipynb")

    print(f"  Notebook 1 : {len(imgs_nb1)} images (cells {list(imgs_nb1.keys())})")
    print(f"  Notebook 2 : {len(imgs_nb2)} images (cells {list(imgs_nb2.keys())})")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    # ── SLIDE 1 : Page de titre ──────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, DARK_BLUE)

    # Ligne décorative
    add_shape_bg(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.05), ACCENT_GREEN)

    add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.2),
                 "Scoring de Crédit et\nPrévision de Défauts",
                 size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.5),
                 "Rapport Intermédiaire — Mars 2026",
                 size=22, color=RGBColor(0xBD, 0xC3, 0xC7), alignment=PP_ALIGN.CENTER)

    add_shape_bg(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.03), MEDIUM_GRAY)

    add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.4),
                 "LP Data Mining — Université Gustave Eiffel",
                 size=16, color=RGBColor(0x95, 0xA5, 0xA6), alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.8),
                 "S. Hamjah  |  I. Naït Daoud  |  M. Trujillo  |  V. Lawrence Raj  |  A. Ngah\nEncadrant : M. Bosco",
                 size=14, color=RGBColor(0x95, 0xA5, 0xA6), alignment=PP_ALIGN.CENTER)

    # ── SLIDE 2 : Sommaire ───────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_title_bar(slide, "Sommaire")

    sections = [
        ("1.", "Contexte et objectifs du projet"),
        ("2.", "Données et exploration"),
        ("3.", "Nettoyage et qualité des données"),
        ("4.", "Feature Engineering"),
        ("5.", "Modélisation — Régression Logistique (Baseline)"),
        ("6.", "Performances du modèle"),
        ("7.", "Analyse temporelle préliminaire"),
        ("8.", "Conclusions et prochaines étapes"),
    ]
    for idx, (num, text) in enumerate(sections):
        y = Inches(1.6) + Inches(idx * 0.6)
        add_text_box(slide, Inches(1.5), y, Inches(0.8), Inches(0.45),
                     num, size=22, color=LIGHT_BLUE, bold=True)
        add_text_box(slide, Inches(2.3), y, Inches(9), Inches(0.45),
                     text, size=20, color=DARK_GRAY)
        if idx < len(sections) - 1:
            add_shape_bg(slide, Inches(2.3), y + Inches(0.5), Inches(8), Inches(0.01),
                         LIGHT_GRAY)

    # ── SLIDE 3 : Contexte et objectifs ──────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_title_bar(slide, "1. Contexte et Objectifs", "Projet tutoré LP Data Mining 2025-2026")

    add_text_box(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(0.4),
                 "Contexte", size=22, color=MEDIUM_BLUE, bold=True)
    add_bullet_list(slide, Inches(0.6), Inches(2.0), Inches(5.5), Inches(2.5), [
        "Secteur bancaire : gestion du risque de crédit",
        "Besoin de modèles prédictifs pour l'octroi de prêts",
        "Données synthétiques réalistes (~130K dossiers clients)",
        "Période couverte : janvier 2021 à décembre 2025",
    ], size=15)

    add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
                 "Deux axes de travail", size=22, color=MEDIUM_BLUE, bold=True)

    # Axe 1 card
    add_shape_bg(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(2.0), LIGHT_BLUE)
    add_text_box(slide, Inches(7.2), Inches(2.2), Inches(5.1), Inches(0.4),
                 "Axe 1 — Scoring de crédit", size=18, color=WHITE, bold=True)
    add_bullet_list(slide, Inches(7.2), Inches(2.7), Inches(5.1), Inches(1.2), [
        "Classification binaire : P(défaut)",
        "Modèles : LogReg, RF, XGBoost",
        "Métrique principale : AUC-ROC",
    ], size=14, color=WHITE)

    # Axe 2 card
    add_shape_bg(slide, Inches(7), Inches(4.3), Inches(5.5), Inches(2.0), ACCENT_ORANGE)
    add_text_box(slide, Inches(7.2), Inches(4.4), Inches(5.1), Inches(0.4),
                 "Axe 2 — Séries temporelles", size=18, color=WHITE, bold=True)
    add_bullet_list(slide, Inches(7.2), Inches(4.9), Inches(5.1), Inches(1.2), [
        "Prévision du volume de défauts (3 mois)",
        "Modèle SARIMA",
        "Métriques : RMSE, MAE, MAPE",
    ], size=14, color=WHITE)

    # Équipe
    add_text_box(slide, Inches(0.6), Inches(5.0), Inches(5.5), Inches(0.4),
                 "Équipe", size=22, color=MEDIUM_BLUE, bold=True)
    team = [
        "Shafin Hamjah — Chef de projet / Data Engineer",
        "Ismaël Naït Daoud — Chef de projet / Data Engineer",
        "Manon Trujillo — Data Scientist",
        "Virginie Lawrence Raj — Data Analyst / Data Engineer",
        "Ange Ngah — Data Scientist (séries temporelles)",
    ]
    add_bullet_list(slide, Inches(0.6), Inches(5.5), Inches(5.5), Inches(1.8), team, size=13)

    # ── SLIDE 4 : Les données ────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_title_bar(slide, "2. Données et Exploration", "Vue d'ensemble du dataset")

    # KPIs
    kpis = [
        ("130 000", "Lignes"),
        ("66", "Colonnes"),
        ("60", "Mois"),
        ("18.6%", "Taux de défaut"),
        ("1 : 4.4", "Ratio déséquilibre"),
    ]
    for i, (val, label) in enumerate(kpis):
        x = Inches(0.6) + Inches(i * 2.5)
        add_metric_card(slide, x, Inches(1.5), Inches(2.2), Inches(1.0), label, val,
                        color=[LIGHT_BLUE, MEDIUM_BLUE, DARK_BLUE, ACCENT_RED, ACCENT_ORANGE][i])

    # Types de variables
    add_text_box(slide, Inches(0.6), Inches(3.0), Inches(6), Inches(0.4),
                 "Catégories de variables", size=20, color=MEDIUM_BLUE, bold=True)

    var_categories = [
        ("Démographiques", "AGE, SEX, MARITAL_STATUS, EDUCATION_LEVEL, REGION..."),
        ("Professionnelles", "EMPLOYMENT_TYPE, EMPLOYMENT_SENIORITY..."),
        ("Financières", "INCOME, EXPENSES, LOAN_AMOUNT, DTI_RATIO, INTEREST_RATE..."),
        ("Historique bancaire", "SAVINGS/CHECKING_BALANCE (M1-M3), VAR_BALANCE_3M..."),
        ("Comportement crédit", "CREDIT_SCORE, NB_OPEN_LOANS, CREDIT_UTILIZATION..."),
        ("Comportementales", "Transactions POS/ATM/Online, connexions, visites..."),
        ("Temporelles", "DATE_MONTH, DEFAULTS_ORIGINATION"),
    ]
    for idx, (cat, desc) in enumerate(var_categories):
        y = Inches(3.5) + Inches(idx * 0.5)
        add_text_box(slide, Inches(0.8), y, Inches(2.8), Inches(0.4),
                     cat, size=13, color=DARK_BLUE, bold=True)
        add_text_box(slide, Inches(3.6), y, Inches(8), Inches(0.4),
                     desc, size=12, color=MEDIUM_GRAY)

    # Image TARGET distribution
    if 12 in imgs_nb1:
        slide.shapes.add_picture(str(imgs_nb1[12]), Inches(7.5), Inches(3.0),
                                  width=Inches(5.5))

    # ── SLIDE 5 : Distribution TARGET + temporelle ──────
    slide = prs.slides.add_slide(blank_layout)
    slide_title_bar(slide, "2. Exploration — Variable cible et tendance temporelle")

    if 12 in imgs_nb1:
        slide.shapes.add_picture(str(imgs_nb1[12]), Inches(0.3), Inches(1.5),
                                  width=Inches(6.2))
    if 23 in imgs_nb1:
        slide.shapes.add_picture(str(imgs_nb1[23]), Inches(6.7), Inches(1.5),
                                  width=Inches(6.2))

    add_text_box(slide, Inches(0.3), Inches(6.8), Inches(12.5), Inches(0.5),
                 "81.4% des clients ne font pas défaut — Classes fortement déséquilibrées (ratio 1:4.4) — Nécessite SMOTE / class_weight",
                 size=13, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 6 : Qualité des données ───────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_title_bar(slide, "3. Nettoyage et Qualité des Données")

    # Valeurs manquantes
    add_text_box(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(0.4),
                 "Valeurs manquantes", size=22, color=MEDIUM_BLUE, bold=True)
    add_bullet_list(slide, Inches(0.6), Inches(2.0), Inches(5.5), Inches(2.0), [
        "1 seule colonne concernée : REGION",
        "11 059 valeurs manquantes (8.51%)",
        "Aucune chaîne vide détectée",
        "0 doublons, 0 CLIENT_ID dupliqués",
    ], size=15)

    # Anomalies
    add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
                 "Anomalies détectées", size=22, color=MEDIUM_BLUE, bold=True)
    add_bullet_list(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(2.0), [
        "DTI > 100% : 10 247 lignes (7.88%)",
        "Soldes négatifs (découverts) :",
        "  Épargne : ~17 500 / mois",
        "  Courant : ~28 800 / mois (min = -2 000)",
        "Aucun âge invalide détecté",
    ], size=15)

    # Résumé nettoyage
    add_text_box(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4),
                 "Résumé du nettoyage", size=22, color=MEDIUM_BLUE, bold=True)

    cleaning_data = [
        ["Étape", "Avant", "Après", "Impact"],
        ["Dimensions", "130 000 × 66", "130 000 × 66", "Aucune suppression"],
        ["Valeurs manquantes", "11 059 (REGION)", "0", "Imputation mode"],
        ["Mémoire", "73.57 MB", "46.37 MB", "Optimisation types"],
        ["Format sauvegarde", "CSV (brut)", "Parquet (12.8 MB)", "Compression efficace"],
    ]
    add_table_slide(slide, Inches(0.6), Inches(5.0), cleaning_data,
                    [Inches(3), Inches(3), Inches(3), Inches(3.5)])

    # ── SLIDE 7 : Feature Engineering ────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_title_bar(slide, "4. Feature Engineering", "38 nouvelles variables créées")

    categories = [
        ("Financières", "11", [
            "SAVINGS_TO_INCOME, LOAN_TO_ANNUAL_INCOME",
            "REPAYMENT_CAPACITY, FINANCIAL_MARGIN",
            "EXPENSE_RATIO, FINANCIAL_STRESS",
        ], LIGHT_BLUE),
        ("Comportementales", "7", [
            "TOTAL_TXN_30D, ONLINE_TXN_RATIO",
            "DIGITAL_ENGAGEMENT, MOBILE_PREFERENCE",
            "WITHDRAWAL_DEPOSIT_RATIO",
        ], ACCENT_GREEN),
        ("Stabilité", "9", [
            "SAVINGS_STABILITY, CHECKING_STABILITY",
            "BALANCE_CV, SAVINGS/CHECKING_TREND",
            "EMPLOYMENT/RESIDENCE_STABILITY",
        ], ACCENT_ORANGE),
        ("Risque", "6", [
            "RISK_SCORE_CUSTOM, CREDIT_SCORE_NORM",
            "LOAN_UTILIZATION, HIGH_RISK_FLAG",
            "IS_NEW_CLIENT",
        ], ACCENT_RED),
        ("Interactions", "5", [
            "AGE_SENIORITY, INCOME_CREDIT_SCORE",
            "DTI_UTILIZATION, LOAN_INTEREST_COST",
            "BALANCE_LOAN_TERM",
        ], MEDIUM_BLUE),
    ]

    for idx, (cat_name, count, features, color) in enumerate(categories):
        x = Inches(0.4) + Inches(idx * 2.55)
        # Header card
        add_shape_bg(slide, x, Inches(1.5), Inches(2.4), Inches(0.8), color)
        add_text_box(slide, x, Inches(1.55), Inches(2.4), Inches(0.4),
                     cat_name, size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(1.95), Inches(2.4), Inches(0.3),
                     f"{count} features", size=12, color=WHITE, alignment=PP_ALIGN.CENTER)
        # Features list
        add_bullet_list(slide, x, Inches(2.5), Inches(2.4), Inches(2.0),
                        features, size=10, color=DARK_GRAY, bullet_char="—")

    # Image corrélation
    if 20 in imgs_nb2:
        slide.shapes.add_picture(str(imgs_nb2[20]), Inches(0.5), Inches(4.5),
                                  width=Inches(12))

    # ── SLIDE 8 : Préparation modélisation ───────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_title_bar(slide, "5. Préparation à la Modélisation",
                    "Split temporel, standardisation, rééquilibrage")

    # Split temporel
    add_shape_bg(slide, Inches(0.6), Inches(1.5), Inches(3.8), Inches(3.0), LIGHT_GRAY)
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(3.4), Inches(0.4),
                 "Split Temporel", size=20, color=DARK_BLUE, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(3.4), Inches(2.2), [
        "Date de coupure : 01/01/2025",
        "Train : 98 700 obs (75.9%)",
        "  Période < 2025-01",
        "  Taux de défaut : 18.69%",
        "Test : 31 300 obs (24.1%)",
        "  Période >= 2025-01",
        "  Taux de défaut : 18.32%",
        "Pas de random split !",
    ], size=14)

    # Standardisation
    add_shape_bg(slide, Inches(4.8), Inches(1.5), Inches(3.8), Inches(3.0), LIGHT_GRAY)
    add_text_box(slide, Inches(5.0), Inches(1.6), Inches(3.4), Inches(0.4),
                 "Standardisation", size=20, color=DARK_BLUE, bold=True)
    add_bullet_list(slide, Inches(5.0), Inches(2.1), Inches(3.4), Inches(2.2), [
        "StandardScaler (scikit-learn)",
        "Fit sur train uniquement",
        "Transform sur train + test",
        "99 features numériques",
        "11 catégorielles encodées (Label)",
    ], size=14)

    # SMOTE
    add_shape_bg(slide, Inches(9.0), Inches(1.5), Inches(3.8), Inches(3.0), LIGHT_GRAY)
    add_text_box(slide, Inches(9.2), Inches(1.6), Inches(3.4), Inches(0.4),
                 "SMOTE (Rééquilibrage)", size=20, color=DARK_BLUE, bold=True)
    add_bullet_list(slide, Inches(9.2), Inches(2.1), Inches(3.4), Inches(2.2), [
        "Avant SMOTE :",
        "  Classe 0 : 80 254",
        "  Classe 1 : 18 446",
        "Après SMOTE :",
        "  Classe 0 : 80 254",
        "  Classe 1 : 80 254",
        "Total : 160 508 observations",
        "Appliqué uniquement sur train",
    ], size=14)

    # Note importante
    add_shape_bg(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(0.8), ACCENT_ORANGE)
    add_text_box(slide, Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.6),
                 "Le split temporel garantit qu'aucune information future ne fuit dans l'entraînement (pas de data leakage). "
                 "Le SMOTE est appliqué UNIQUEMENT sur le jeu d'entraînement pour ne pas biaiser l'évaluation.",
                 size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 9 : Résultats du modèle ───────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_title_bar(slide, "6. Régression Logistique — Performances",
                    "Modèle baseline avec class_weight='balanced'")

    # Métriques cards
    metrics = [
        ("AUC-ROC", "0.7105", LIGHT_BLUE),
        ("Accuracy", "69.9%", MEDIUM_BLUE),
        ("Precision", "32.2%", ACCENT_ORANGE),
        ("Recall", "58.0%", ACCENT_GREEN),
        ("F1-Score", "0.414", DARK_BLUE),
    ]
    for i, (label, value, color) in enumerate(metrics):
        x = Inches(0.6) + Inches(i * 2.5)
        add_metric_card(slide, x, Inches(1.5), Inches(2.2), Inches(1.1), label, value, color)

    # Classification report table
    report_data = [
        ["Classe", "Precision", "Recall", "F1-Score", "Support"],
        ["Non défaut (0)", "0.89", "0.73", "0.80", "25 566"],
        ["Défaut (1)", "0.32", "0.58", "0.41", "5 734"],
        ["Macro avg", "0.60", "0.65", "0.61", "31 300"],
        ["Weighted avg", "0.78", "0.70", "0.73", "31 300"],
    ]
    add_table_slide(slide, Inches(0.6), Inches(3.0), report_data,
                    [Inches(2.5), Inches(2), Inches(2), Inches(2), Inches(2)])

    # CV
    add_text_box(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
                 "Validation croisée (5-fold) : AUC-ROC = 0.717 ± 0.011",
                 size=16, color=DARK_BLUE, bold=True)
    add_text_box(slide, Inches(0.6), Inches(5.6), Inches(12), Inches(0.5),
                 "Folds : 0.716 | 0.716 | 0.717 | 0.710 | 0.727 — Modèle stable, pas de surapprentissage",
                 size=13, color=MEDIUM_GRAY)

    # Interprétation
    add_shape_bg(slide, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.8), LIGHT_GRAY)
    add_bullet_list(slide, Inches(0.8), Inches(6.25), Inches(11.7), Inches(0.7), [
        "Recall 58% : le modèle détecte plus de la moitié des défauts — bon pour un baseline",
        "Precision 32% : 1 alerte sur 3 est un vrai défaut — à améliorer avec des modèles avancés",
    ], size=13, color=DARK_GRAY)

    # ── SLIDE 10 : Graphiques de performance ─────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_title_bar(slide, "6. Visualisations des Performances",
                    "Courbe ROC, Matrice de confusion, Precision-Recall, Distribution des scores")

    if 38 in imgs_nb2:
        slide.shapes.add_picture(str(imgs_nb2[38]), Inches(0.5), Inches(1.4),
                                  width=Inches(12.3))

    # ── SLIDE 11 : Importance des features ───────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_title_bar(slide, "6. Importance des Variables",
                    "Top 15 features — Coefficients de la régression logistique")

    # Table top features
    top_features = [
        ["Rang", "Variable", "Coefficient", "Effet"],
        ["1", "DTI_RATIO", "+0.427", "Augmente le risque"],
        ["2", "INSTALLMENT_TO_INCOME", "+0.427", "Augmente le risque"],
        ["3", "VAR_BALANCE_3M", "-0.220", "Diminue le risque"],
        ["4", "CREDIT_SCORE_NORM", "-0.172", "Diminue le risque"],
        ["5", "CREDIT_SCORE", "-0.172", "Diminue le risque"],
        ["6", "CHECKING_STABILITY", "+0.129", "Augmente le risque"],
        ["7", "HAS_PREVIOUS_DEFAULT", "+0.117", "Augmente le risque"],
        ["8", "LOAN_AMOUNT", "+0.110", "Augmente le risque"],
        ["9", "INSURANCE_TAKEN", "-0.108", "Diminue le risque"],
        ["10", "PHONE_VERIFIED", "+0.089", "Augmente le risque"],
    ]
    add_table_slide(slide, Inches(0.6), Inches(1.5), top_features,
                    [Inches(1), Inches(3.5), Inches(2), Inches(3)])

    if 41 in imgs_nb2:
        slide.shapes.add_picture(str(imgs_nb2[41]), Inches(6.8), Inches(1.5),
                                  width=Inches(6))

    add_text_box(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 "Le ratio d'endettement (DTI) et le rapport mensualité/revenu sont les prédicteurs les plus forts du défaut.",
                 size=14, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 12 : Analyse temporelle ────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_title_bar(slide, "7. Analyse Temporelle Préliminaire",
                    "Série DEFAULTS_ORIGINATION — 60 mois")

    # Stats
    add_text_box(slide, Inches(0.6), Inches(1.5), Inches(4.5), Inches(0.4),
                 "Statistiques de la série", size=20, color=MEDIUM_BLUE, bold=True)

    ts_stats = [
        ["Statistique", "Valeur"],
        ["Moyenne", "403 défauts/mois"],
        ["Écart-type", "98.3"],
        ["Minimum", "226"],
        ["Maximum", "653"],
        ["Médiane", "405"],
        ["Tendance", "CROISSANTE (r = 0.583)"],
    ]
    add_table_slide(slide, Inches(0.6), Inches(2.0), ts_stats,
                    [Inches(2.5), Inches(3)])

    # Saisonnalité
    add_text_box(slide, Inches(0.6), Inches(5.0), Inches(4.5), Inches(0.4),
                 "Pics saisonniers", size=18, color=MEDIUM_BLUE, bold=True)
    add_bullet_list(slide, Inches(0.6), Inches(5.4), Inches(4.5), Inches(1.5), [
        "Décembre : 525 défauts (pic annuel)",
        "Septembre : 500 défauts",
        "Mars : 484 défauts",
    ], size=14)

    if 45 in imgs_nb2:
        slide.shapes.add_picture(str(imgs_nb2[45]), Inches(6.3), Inches(1.4),
                                  width=Inches(6.5))

    # ── SLIDE 13 : Saisonnalité ──────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_title_bar(slide, "7. Analyse de la Saisonnalité",
                    "Moyenne des défauts par mois de l'année")

    if 47 in imgs_nb2:
        slide.shapes.add_picture(str(imgs_nb2[47]), Inches(1.5), Inches(1.4),
                                  width=Inches(10))

    add_text_box(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 "Saisonnalité marquée : pics en fin d'année (décembre) et en fin de trimestre — "
                 "à intégrer dans le modèle SARIMA (composante saisonnière S=12)",
                 size=14, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 14 : Conclusions ───────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_title_bar(slide, "8. Conclusions et Prochaines Étapes")

    # Bilan
    add_text_box(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.4),
                 "Bilan Phase 1", size=22, color=MEDIUM_BLUE, bold=True)

    add_shape_bg(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.5), LIGHT_GRAY)
    bilan = [
        "Exploration complète : 130K lignes, 66 colonnes",
        "Nettoyage : 1 variable avec manquants traitée",
        "Feature engineering : 38 nouvelles features",
        "5 catégories : financières, comportementales, stabilité, risque, interactions",
        "Modèle baseline (LogReg) : AUC-ROC = 0.71",
        "Validation croisée stable : 0.717 ± 0.011",
        "Analyse temporelle : tendance croissante + saisonnalité",
        "Tous les résultats sauvegardés (parquet, joblib, JSON)",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.4), Inches(4.3),
                    bilan, size=14, color=DARK_GRAY, bullet_char="\u2713")

    # Prochaines étapes
    add_text_box(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(0.4),
                 "Phase 2 — Prochaines étapes", size=22, color=ACCENT_RED, bold=True)

    add_shape_bg(slide, Inches(7), Inches(2.0), Inches(5.8), Inches(4.5), LIGHT_GRAY)
    next_steps = [
        "Optimisation hyperparamètres (GridSearch/Bayesian)",
        "Modèles avancés : Random Forest, XGBoost, LightGBM",
        "Sélection de variables (importance, RFE)",
        "Calibration des probabilités",
        "Modèle SARIMA(p,d,q)(P,D,Q,12)",
        "Tests de stationnarité (ADF, KPSS)",
        "Prévisions 3 mois avec intervalles de confiance",
        "Dashboard de visualisation",
    ]
    add_bullet_list(slide, Inches(7.2), Inches(2.1), Inches(5.4), Inches(4.3),
                    next_steps, size=14, color=DARK_GRAY, bullet_char="\u27A4")

    # Timeline
    add_shape_bg(slide, Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.5), DARK_BLUE)
    add_text_box(slide, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.4),
                 "Phase 2 : Janvier — Mars 2026  |  Phase 3 (finale) : Avril — Juin 2026  |  Soutenance : Juin 2026",
                 size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 15 : Merci ─────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, DARK_BLUE)

    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.0),
                 "Merci de votre attention", size=40, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_shape_bg(slide, Inches(4), Inches(3.8), Inches(5.3), Inches(0.04), ACCENT_GREEN)

    add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.5),
                 "Questions ?", size=28, color=RGBColor(0xBD, 0xC3, 0xC7),
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.8),
                 "Scoring de Crédit et Prévision de Défauts\n"
                 "LP Data Mining — Université Gustave Eiffel — 2025-2026",
                 size=16, color=RGBColor(0x95, 0xA5, 0xA6), alignment=PP_ALIGN.CENTER)

    # ── Sauvegarde ───────────────────────────────────────
    output_path = ROOT / "outputs" / "reports" / "presentation_scoring_credit.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"\nPresentation sauvegardee : {output_path}")
    print(f"  {len(prs.slides)} slides generees")


if __name__ == "__main__":
    main()
